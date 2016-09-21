import csv
import docx
import exifread
import json
import logging
import os
import os.path
import magic
import shapefile
import subprocess
import xlrd
import xmltodict
import zipfile
from dbfread       import DBF
from elasticsearch import Elasticsearch
from openpyxl      import load_workbook
from pptx          import Presentation
from megadex.util  import get_ext


logger = logging.getLogger('megadex')


class Indexer(Elasticsearch):

    def __init__(self, *args, **kwargs):
        self.deep      = kwargs.pop('deep', False)
        self.paged_pdf = kwargs.pop('paged_pdf', False)
        super(Indexer, self).__init__(*args, **kwargs)

    def index_file(self, filename):
        ext = get_ext(filename)
        file_type = magic.from_file(filename)
        mime      = magic.from_file(filename, mime=True)
        logger.info('indexing: {0} {1}'.format(filename, mime))

        # first try to index purely off extension
        try:
            cb  = getattr(self, 'index_{0}'.format(ext.replace('.', '')))
            res = cb(filename)
            if res is not None:
                logger.info(res)
            return
        except Exception as e:
            logger.error(e)

        # last ditch index based on magic file type
        try:
            logger.info(self._index_magic(filename))
        except Exception as e:
            logger.error(e)

    def _index_magic(self, filename):
        file_type = magic.from_file(filename).lower()
        mime      = magic.from_file(filename, mime=True).lower()
        logger.debug('magic file type: {0} {1}'.format(file_type, mime))

	if 'access' in (file_type, mime):
	    return self.index_mdb(filename)

        if 'ascii' in (file_type, mime):
            return self.index_text(filename)

        if 'csv' in (file_type, mime):
            return self.index_csv(filename)

        if 'jpg' in (file_type, mime):
            return self.index_jpg(filename)

        if 'gif' in (file_type, mime):
            return self.index_jpg(filename)

        if 'png' in (file_type, mime):
            return self.index_png(filename)

        if 'pdf' in (file_type, mime):
            return self.index_pdf(filename)

        if 'tiff' in (file_type, mime):
            return self.index_tiff(filename)

        if 'xml' in (file_type, mime):
            return self.index_xml(filename)

        if 'zip' in (file_type, mime):
            return self.index_zip(filename)

        if 'text' in (file_type, mime):
            return self.index_text(filename)

    def index_asc(self, filename):
        with open(filename) as f:
            meta = str(f.read())
            data = {'text': meta}

        data['filename'] = filename

        return self.index(index="asc", doc_type="asc", body=data)

    def index_csv(self, filename):
        with open(filename, mode='r') as infile:
            reader = csv.reader(infile)
            data = {
                rows[0].replace('.', ''):rows[1]
                for rows in reader
            }
            data['filename'] = filename

        return self.index(index="csv", doc_type="csv", body=data)

    def index_doc(self, filename):
        document = docx.Document(filename)
        data = '\n'.join([
            paragraph.text.encode('utf-8')
            for paragraph in document.paragraphs
        ])
        data['filename'] = filename

        return self.index(index="doc", doc_type="doc", body=data)

    def index_docx(self, filename):
        document = docx.Document(filename)
        data = {}
        data['paragraphs'] = [
            paragraph.text.encode('utf-8')
            for paragraph in document.paragraphs
        ]

        data['filename'] = filename

        return self.index(index="docx", doc_type="docx", body=data)

    def index_jpg(self, filename):
        f = open(filename, 'rb')
        data = dict( (k, str(v)) for k, v in exifread.process_file(f).items())
        data['filename'] = filename

        return self.index(index="jpg", doc_type="jpg", body=data)

    def index_gif(self, filename):
        f = open(filename, 'rb')
        data = dict( (k, str(v)) for k, v in exifread.process_file(f).items())
        data['filename'] = filename

        return self.index(index="gif", doc_type="gif", body=data)

    def index_kml(self, filename):
        data = {}

        # XXX: make this smart
        if os.path.getsize(filename) < 1000000:
            data = xmltodict.parse(open(filename).read())

        data['filename'] = filename

        return self.index(index="kml", doc_type="kml", body=data)

    def index_kmz(self, filename):
        if not zipfile.is_zipfile(filename):
            raise TypeError("{0} is not a zipfile".format(filename))

        data = {'files':[]}

        zf = zipfile.ZipFile(filename)

        for info in zf.infolist():
            data['files'].append({
                'filename':        info.filename,
                'comment':         info.comment,
                'version':         info.create_version,
                'compressed_size': info.compress_size,
                'size':            info.file_size,
            })

        data['filename'] = filename

        return self.index(index="kmz", doc_type="kmz", body=data)

    def index_mdb(self, filename):
        tables = subprocess.Popen(
	    ["mdb-tables", "-d,", "-t", "table", filename],
            stdout=subprocess.PIPE
        ).communicate()[0]

	tables = tables.rstrip('\n').split(',')

	data = {}
	data['tables'] = []

	for table in tables:
	    raw_data = subprocess.Popen(
		["mdb-export", filename, table],
		stdout=subprocess.PIPE
	    ).communicate()[0]

	    raw_data = raw_data.rstrip('\n').split('\n')

	    if not filter(None, raw_data):
		continue

	    d = {}
	    d['table'] = table
	    d['data'] = {}

	    headers = []
            for line in raw_data:
		if not headers:
		    headers = line.split(',')
		    for header in headers:
		        d['data'][header] = []
		    continue

		for col, v in zip(headers, line.split(',')):
		    d['data'][col].append(v.replace('"', '').replace("'", ""))

	    d['columns'] = headers
	    data['tables'].append(d)

	if not data:
	    return

	data['filename'] = filename

	return self.index(index="mdb", doc_type="mdb", body=data)

    def index_dbf(self, filename):
	try:
	    table = DBF(
		filename,
		ignorecase              = True,
		ignore_missing_memofile = True,
	    )
	except Exception as e:
	    try:
		table = DBF(
		    filename,
		    ignorecase              = True,
		    ignore_missing_memofile = True,
		    raw                     = True,
		    encoding                = 'iso-8859-1',
		)
	    except Exception as ee:
		table = DBF(
		    filename,
		    ignorecase              = True,
		    ignore_missing_memofile = True,
		    encoding                = 'utf8',
		)

	data = {}
	data['fields']     = table.field_names
	data['db_version'] = table.header.dbversion
	data['year']       = table.header.year
	data['month']      = table.header.month
	data['day']        = table.header.day
	data['records']    = table.header.numrecords

	data['filename'] = filename

	# XXX: get more data
	# https://dbfread.readthedocs.io/en/latest/dbf_objects.html

        return self.index(index="dbf", doc_type="dbf", body=data)

    def index_pdf(self, filename):
        output = subprocess.Popen(
            ["pdftotext", "-q", "-eol", "unix", filename, "-"],
            stdout=subprocess.PIPE
        ).communicate()[0]
        data = {'text': output}

        if self.paged_pdf:
            trunc_filename = filename.rstrip('.pdf')
            try:
                page = trunc_filename[trunc_filename.rindex('-')+1:]
            except Exception as e:
                logger.warn(e.info)
                page = '1'

            data['page'] = page
            data['filename'] = trunc_filename.rstrip('-' + page) + '.pdf'
        else:
            data['filename'] = filename

        return self.index(index="pdf", doc_type="pdf", body=data)

    def index_png(self, filename):
        f = open(filename, 'rb')
        data = exifread.process_file(f)
        data['filename'] = filename

        return self.index(index="png", doc_type="png", body=data)

    def index_ppt(self, filename):
	# XXX:
	pass

    def index_pptx(self, filename):
	prs = Presentation(filename)

        data = {}
	# text_runs will be populated with a list of strings,
	# one for each text run in presentation

        data['filename'] = filename

        res = None

	for i, slide in enumerate(prs.slides):
	    for shape in slide.shapes:
		if not shape.has_text_frame:
		    continue

                text_runs = []
		for paragraph in shape.text_frame.paragraphs:
		    for run in paragraph.runs:
			text_runs.append(run.text)

                d = data
                d['page'] = i + 1
                d['text'] = "\n".join(text_runs)

                res = self.index(index="pptx", doc_type="pptx", body=d)

        return res

    def index_scn(self, filename):
        with open(filename) as f:
            meta = str(f.read())
            data = {'text': meta}

        data['filename'] = filename

        return self.index(index="scn", doc_type="scn", body=data)

    def index_shp(self, filename):
        sf = shapefile.Reader(filename)
        shapes = sf.shapes()
        logger.info("logging: {0}".format(shapes))

    def index_tiff(self, filename):
        f = open(filename, 'rb')
        data = dict( (k, str(v)) for k, v in exifread.process_file(f).items())
        data['filename'] = filename

        return self.index(index="tiff", doc_type="tiff", body=data)

    def index_txt(self, filename):
        with open(filename) as f:
            meta = str(f.read())
            data = {'text': meta}

        data['filename'] = filename

        return self.index(index="txt", doc_type="txt", body=data)

    def index_xls(self, filename):
        book = xlrd.open_workbook(filename)
        sheets = book.sheet_names()
        data = {}
        data['sheets'] = []

        for sheet_name in sheets:
            d = {}
            sheet = book.sheet_by_name(sheet_name)
            d['sheet'] = sheet
            for rx in range(sheet.nrows):
                # XXX
                pass

            data['sheets'].append(d)

        data['filename'] = filename

        return self.index(index="xls", doc_type="xls", body=data)

    def index_xlsx(self, filename):
        wb = load_workbook(filename, use_iterators=True)
        sheets = wb.get_sheet_names()
        logger.info("logging: {0}".format(sheets))
        data = {}
        data['sheets'] = []
        for sheet_name in sheets:
            sheet   = wb[sheet_name]
            headers = []
            d       = {}
            first   = True
            for row in sheet.iter_rows():
                if first:
                    headers = [ x.value for x in row ]
                    for x in row:
                        d[x.value] = []
                    first = False
                    continue
                else:
                    for col, x in zip(headers, row):
                        d[col].append(x.value)

            data['sheets'].append(d)

        data['filename'] = filename

        return self.index(index="xlsx", doc_type="xlsx", body=data)

    def index_xml(self, filename):
        data = xmltodict.parse(open(filename).read())

        data['filename'] = filename

        return self.index(index="xml", doc_type="xml", body=data)

    def index_xyz(self, filename):
        return self._index_magic(filename)

    def index_zip(self, filename):
        if not zipfile.is_zipfile(filename):
            raise TypeError("{0} is not a zipfile".format(filename))

        data = {'files':[]}

        zf = zipfile.ZipFile(filename)

        for info in zf.infolist():
            d = {
                'filename':        info.filename,
                'comment':         info.comment,
                'version':         info.create_version,
                'compressed_size': info.compress_size,
                'size':            info.file_size,
            }

            if self.deep:
                # XXX: handle this
                pass

            data['files'].append(d)

        data['filename'] = filename

        return self.index(index="zip", doc_type="zip", body=data)
